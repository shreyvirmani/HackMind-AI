from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE


def export_ppt(
    deck,
    research=None,
    judge=None,
    include_ai_slides=False,
):

    prs = Presentation()

    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for index, slide_data in enumerate(deck.slides):

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(245, 247, 250)

        # Title Banner
        banner = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            0,
            0,
            prs.slide_width,
            Inches(0.8),
        )

        banner.fill.solid()
        banner.fill.fore_color.rgb = RGBColor(37, 99, 235)
        banner.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.4),
            Inches(0.15),
            Inches(10),
            Inches(0.5),
        )

        title = title_box.text_frame
        p = title.paragraphs[0]
        p.text = slide_data.title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        body_box = slide.shapes.add_textbox(
            Inches(0.8),
            Inches(1.2),
            Inches(11.5),
            Inches(5.5),
        )

        tf = body_box.text_frame

        for point in slide_data.content:

            para = tf.add_paragraph()

            para.text = "• " + point
            para.level = 0
            para.font.size = Pt(22)
            para.font.color.rgb = RGBColor(40, 40, 40)

        # Footer
        footer = slide.shapes.add_textbox(
            Inches(0.4),
            Inches(6.9),
            Inches(12),
            Inches(0.3),
        )

        f = footer.text_frame.paragraphs[0]
        f.text = f"HackMind AI | Slide {index+1}/{len(deck.slides)}"
        f.font.size = Pt(10)
        f.font.color.rgb = RGBColor(120, 120, 120)
        f.alignment = PP_ALIGN.RIGHT
    
    # ==========================
    # AI APPENDIX
    # ==========================

    if include_ai_slides:

        if judge:

            add_judge_slide(prs, judge)

        if research:

            add_research_slide(prs, research)

    buffer = BytesIO()

    prs.save(buffer)

    buffer.seek(0)

    return buffer

def add_judge_slide(prs, judge):

    slide = prs.slides.add_slide(prs.slide_layouts[1])

    slide.shapes.title.text = "🏆 AI Judge Evaluation"

    tf = slide.placeholders[1].text_frame

    tf.text = f"Overall Score: {judge.overall_score}/100"

    for strength in judge.strengths:

        p = tf.add_paragraph()

        p.text = f"✓ {strength}"

def add_research_slide(prs, research):

    slide = prs.slides.add_slide(prs.slide_layouts[1])

    slide.shapes.title.text = "🔍 AI Research Insights"

    tf = slide.placeholders[1].text_frame

    tf.text = "Top Competitors"

    for competitor in research.competitors:

        p = tf.add_paragraph()

        p.text = competitor.name


